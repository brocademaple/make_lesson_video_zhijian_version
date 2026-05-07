"""根据 CourseSlide.type 渲染单页（python-pptx）。"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Callable, Dict

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

import sys

_ROOT = Path(__file__).resolve().parent.parent
_TEMPL = _ROOT / "templates"
if str(_TEMPL) not in sys.path:
    sys.path.insert(0, str(_TEMPL))

import theme  # type: ignore

from models import CourseSlide

logger = logging.getLogger(__name__)

C = theme.COLORS


def _font_title(paragraph, size: int = 32, bold: bool = True, color=None):
    paragraph.font.name = theme.FONT_CN
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color or C["primary"]


def _font_body(paragraph, size: int = 18, color=None):
    paragraph.font.name = theme.FONT_CN
    paragraph.font.size = Pt(size)
    paragraph.font.bold = False
    paragraph.font.color.rgb = color or C["text"]


def _add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _load_picture(slide, asset_path: Path, left, top, width, height):
    if asset_path.is_file():
        try:
            slide.shapes.add_picture(str(asset_path), left, top, width=width, height=height)
            return True
        except Exception as e:
            logger.debug("插图失败 %s: %s", asset_path, e)
    return False


def render_title(prs: Presentation, cs: CourseSlide, assets: Path) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # 背景条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), theme.SLIDE_WIDTH, Inches(2.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C["primary"]
    shape.line.fill.background()
    tf = _add_textbox(slide, Inches(0.6), Inches(0.55), Inches(12), Inches(1.1)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = cs.title
    p.alignment = PP_ALIGN.LEFT
    _font_title(p, 36, True, C["white"])

    if cs.subtitle:
        p2 = tf.add_paragraph()
        p2.text = cs.subtitle
        p2.space_before = Pt(6)
        _font_body(p2, 20, C["white"])

    tag = "新兵营 · 质检培训"
    p3 = tf.add_paragraph()
    p3.text = tag
    p3.space_before = Pt(14)
    _font_body(p3, 14, RGBColor(0xDD, 0xE4, 0xFF))

    _load_picture(slide, assets / "shield.png", Inches(11), Inches(0.35), Inches(1.4), Inches(1.4))


def render_agenda(prs: Presentation, cs: CourseSlide, assets: Path) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tf = _add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(1)).text_frame
    p = tf.paragraphs[0]
    p.text = cs.title
    _font_title(p, 30)

    body = _add_textbox(slide, Inches(0.7), Inches(1.35), Inches(11.5), Inches(5.5)).text_frame
    body.word_wrap = True
    items = (cs.bullets or [])[:3]
    if not items and cs.main_text:
        items = [cs.main_text]
    for i, line in enumerate(items):
        para = body.paragraphs[0] if i == 0 else body.add_paragraph()
        prefix = f"{i+1}. "
        para.text = prefix + line
        para.space_after = Pt(12)
        _font_body(para, 22)


def render_transition(prs: Presentation, cs: CourseSlide, assets: Path) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    mid = _add_textbox(slide, Inches(1), Inches(2.4), Inches(11), Inches(2)).text_frame
    mid.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = mid.paragraphs[0]
    p.text = cs.title
    p.alignment = PP_ALIGN.CENTER
    _font_title(p, 34, True, C["accent_purple"])
    if cs.main_text:
        p2 = mid.add_paragraph()
        p2.text = cs.main_text
        p2.space_before = Pt(16)
        p2.alignment = PP_ALIGN.CENTER
        _font_body(p2, 20, C["muted"])


def render_rule_card(prs: Presentation, cs: CourseSlide, assets: Path) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    left_x = Inches(0.55)
    top_y = Inches(0.45)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left_x,
        top_y,
        Inches(12.2),
        Inches(6.2),
    )
    try:
        card.adjustments[0] = 0.08
    except Exception:
        pass
    card.fill.solid()
    card.fill.fore_color.rgb = C["card_bg"]
    card.line.color.rgb = C["primary"]

    icon_x = left_x + Inches(0.35)
    _load_picture(slide, assets / "shield.png", icon_x, top_y + Inches(0.35), Inches(0.75), Inches(0.75))

    tf = _add_textbox(slide, left_x + Inches(1.15), top_y + Inches(0.35), Inches(10.5), Inches(0.9)).text_frame
    p = tf.paragraphs[0]
    p.text = cs.title
    _font_title(p, 26)

    if cs.main_text:
        t2 = _add_textbox(slide, left_x + Inches(0.45), top_y + Inches(1.25), Inches(11.5), Inches(1)).text_frame
        p2 = t2.paragraphs[0]
        p2.text = cs.main_text
        _font_body(p2, 20, C["alert_red"])

    bx = _add_textbox(slide, left_x + Inches(0.45), top_y + Inches(2.45), Inches(11.5), Inches(3.6)).text_frame
    bx.word_wrap = True
    bullets = (cs.bullets or [])[:3]
    for i, b in enumerate(bullets):
        para = bx.paragraphs[0] if i == 0 else bx.add_paragraph()
        if i > 0:
            para = bx.add_paragraph()
        para.text = "• " + b
        para.space_after = Pt(10)
        _font_body(para, 17)

    _load_picture(slide, assets / "warning.png", Inches(11.5), top_y + Inches(4.9), Inches(1.1), Inches(1.1))


def render_case_dialogue(prs: Presentation, cs: CourseSlide, assets: Path) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tf = _add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12), Inches(0.8)).text_frame
    p = tf.paragraphs[0]
    p.text = cs.title
    _font_title(p, 26)

    y = 1.25
    for i, row in enumerate(cs.dialogue[:6]):
        speaker = str(row.get("speaker", "角色"))
        text = str(row.get("text", ""))
        is_user = "用户" in speaker or speaker.lower() == "user"
        bg = RGBColor(0xE8, 0xEE, 0xFF) if is_user else RGBColor(0xFF, 0xF3, 0xE0)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.6 if is_user else 4.2),
            Inches(y),
            Inches(7.8 if is_user else 8.5),
            Inches(1.05),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = bg
        box.line.fill.background()
        bt = box.text_frame
        bt.word_wrap = True
        pr = bt.paragraphs[0]
        pr.text = f"{speaker}：{text}"
        _font_body(pr, 15)
        y += 1.15

    risk = cs.main_text or "风险：请在实际业务中守住红线"
    rf = _add_textbox(slide, Inches(0.55), Inches(6.35), Inches(12), Inches(0.75)).text_frame
    pr = rf.paragraphs[0]
    pr.text = risk
    _font_body(pr, 14, C["alert_red"])

    _load_picture(slide, assets / "cross.png", Inches(11.8), Inches(1.2), Inches(0.9), Inches(0.9))
    _load_picture(slide, assets / "check.png", Inches(11.8), Inches(4.0), Inches(0.9), Inches(0.9))


def render_quiz(prs: Presentation, cs: CourseSlide, assets: Path) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tf = _add_textbox(slide, Inches(0.55), Inches(0.4), Inches(12), Inches(1.2)).text_frame
    p = tf.paragraphs[0]
    qtext = ""
    opts: list[str] = []
    if cs.quiz and isinstance(cs.quiz, dict):
        qtext = str(cs.quiz.get("question", cs.title))
        opts = list(cs.quiz.get("options") or [])[:4]
    else:
        qtext = cs.title
    p.text = qtext
    _font_title(p, 22)

    hint = _add_textbox(slide, Inches(0.55), Inches(1.45), Inches(12), Inches(0.5)).text_frame
    hp = hint.paragraphs[0]
    hp.text = "先想 3 秒，再对照解析。"
    _font_body(hp, 14, C["muted"])

    body = _add_textbox(slide, Inches(0.75), Inches(2.15), Inches(11.5), Inches(4.5)).text_frame
    for i, opt in enumerate(opts):
        para = body.paragraphs[0] if i == 0 else body.add_paragraph()
        if i > 0:
            para = body.add_paragraph()
        label = chr(65 + i)
        para.text = f"{label}. {opt}"
        para.space_after = Pt(12)
        _font_body(para, 18)

    _load_picture(slide, assets / "phone.png", Inches(11.2), Inches(5.6), Inches(1.3), Inches(1.3))


def render_explanation(prs: Presentation, cs: CourseSlide, assets: Path) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tf = _add_textbox(slide, Inches(0.55), Inches(0.4), Inches(12), Inches(0.9)).text_frame
    p = tf.paragraphs[0]
    p.text = cs.title or "解析"
    _font_title(p, 26)

    ans = cs.main_text or "正确答案见原文规则。"
    t2 = _add_textbox(slide, Inches(0.55), Inches(1.25), Inches(12), Inches(1)).text_frame
    p2 = t2.paragraphs[0]
    p2.text = "答案：" + ans
    _font_body(p2, 20, RGBColor(0x1B, 0x7E, 0x4B))

    exp = cs.explanation or cs.narration
    t3 = _add_textbox(slide, Inches(0.55), Inches(2.45), Inches(12), Inches(4)).text_frame
    t3.word_wrap = True
    p3 = t3.paragraphs[0]
    p3.text = exp
    _font_body(p3, 17)

    _load_picture(slide, assets / "check.png", Inches(11.5), Inches(0.35), Inches(1.1), Inches(1.1))


def render_summary(prs: Presentation, cs: CourseSlide, assets: Path) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tf = _add_textbox(slide, Inches(0.55), Inches(0.4), Inches(12), Inches(0.9)).text_frame
    p = tf.paragraphs[0]
    p.text = cs.title
    _font_title(p, 28)

    body = _add_textbox(slide, Inches(0.65), Inches(1.35), Inches(11.5), Inches(5.3)).text_frame
    body.word_wrap = True
    lines = list(cs.bullets or [])
    if cs.main_text:
        lines.insert(0, cs.main_text)
    if not lines:
        lines = [cs.narration[:120] + "…"] if cs.narration else ["本节重点回顾"]
    for i, line in enumerate(lines[:5]):
        para = body.paragraphs[0] if i == 0 else body.add_paragraph()
        if i == 0:
            para.text = "本节要点：\n" + line
            _font_body(para, 20, C["primary"])
        else:
            para.text = "✧ " + line
            _font_body(para, 18, C["text"])
        para.space_after = Pt(10)

    foot = _add_textbox(slide, Inches(0.55), Inches(6.6), Inches(12), Inches(0.6)).text_frame
    fp = foot.paragraphs[0]
    fp.text = "下一步：回到业务场景自查三条红线。"
    _font_body(fp, 14, C["muted"])


def render_fallback_unknown(prs: Presentation, cs: CourseSlide, assets: Path) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tf = _add_textbox(slide, Inches(0.55), Inches(2.4), Inches(12), Inches(2)).text_frame
    p = tf.paragraphs[0]
    p.text = cs.title or "内容页"
    _font_title(p, 24)
    p2 = tf.add_paragraph()
    p2.text = cs.main_text or (cs.narration[:200] if cs.narration else "")
    p2.space_before = Pt(12)
    _font_body(p2, 16)


def render_error_slide(prs: Presentation, slide_id: str, err: str) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tf = _add_textbox(slide, Inches(0.55), Inches(2.5), Inches(12), Inches(2.5)).text_frame
    p = tf.paragraphs[0]
    p.text = "页面生成失败"
    _font_title(p, 28, True, C["alert_red"])
    p2 = tf.add_paragraph()
    p2.text = f"{slide_id}\n{err[:500]}"
    p2.space_before = Pt(10)
    _font_body(p2, 14)


RENDERERS: Dict[str, Callable[..., None]] = {
    "title": render_title,
    "agenda": render_agenda,
    "transition": render_transition,
    "rule_card": render_rule_card,
    "case_dialogue": render_case_dialogue,
    "quiz": render_quiz,
    "explanation": render_explanation,
    "summary": render_summary,
}


def render_course_slide(prs: Presentation, cs: CourseSlide, assets_dir: Path) -> None:
    fn = RENDERERS.get(cs.type)
    if fn:
        fn(prs, cs, assets_dir)
    else:
        render_fallback_unknown(prs, cs, assets_dir)
