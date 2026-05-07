"""读取原始 PPTX，抽取每页文本与结构信息。"""

from __future__ import annotations

import logging
import re
from typing import Iterable, List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from models import RawSlide

logger = logging.getLogger(__name__)


def _iter_shapes(shape: BaseShape) -> Iterable[BaseShape]:
    yield shape
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:  # type: ignore[attr-defined]
            yield from _iter_shapes(child)


def _shape_text(shape: BaseShape) -> str:
    if not shape.has_text_frame:
        return ""
    lines: List[str] = []
    for para in shape.text_frame.paragraphs:
        t = "".join(run.text for run in para.runs).strip()
        if t:
            lines.append(t)
    return "\n".join(lines)


def _count_images(shape: BaseShape) -> int:
    n = 0
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return 1
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:  # type: ignore[attr-defined]
            n += _count_images(child)
    return n


def _table_count_on_shape(shape: BaseShape) -> int:
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        return 1
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return sum(_table_count_on_shape(c) for c in shape.shapes)  # type: ignore[attr-defined]
    return 0


def _slide_notes(slide: Slide) -> str | None:
    try:
        if getattr(slide, "has_notes_slide", False):
            tf = slide.notes_slide.notes_text_frame
            t = tf.text.strip()
            return t or None
    except Exception as e:
        logger.debug("读取备注失败: %s", e)
    return None


def _guess_title(slide: Slide, blocks: List[str], all_text: str) -> str | None:
    if slide.shapes.title and slide.shapes.title.text.strip():
        return _collapse_ws(slide.shapes.title.text)
    if blocks:
        first = blocks[0].strip()
        if len(first) <= 80:
            return first
    first_line = all_text.split("\n", 1)[0].strip()
    return _collapse_ws(first_line)[:80] if first_line else None


def _collapse_ws(s: str) -> str:
    return re.sub(r"\s+", " ", s).strip()


def read_presentation(path: str) -> List[RawSlide]:
    prs = Presentation(path)
    raw_slides: List[RawSlide] = []

    for idx, slide in enumerate(prs.slides):
        blocks: List[str] = []
        img_count = 0
        tbl_count = 0
        layout_name = ""
        try:
            layout_name = slide.slide_layout.name or ""
        except Exception:
            pass

        for shape in slide.shapes:
            for s in _iter_shapes(shape):
                t = _shape_text(s)
                if t:
                    blocks.append(t)
                img_count += _count_images(s)
                tbl_count += _table_count_on_shape(s)

        # 去重连续重复块
        dedup: List[str] = []
        for b in blocks:
            if dedup and dedup[-1] == b:
                continue
            dedup.append(b)

        all_text = "\n\n".join(dedup)
        title_text = _guess_title(slide, dedup, all_text)

        notes = _slide_notes(slide)
        layout_info = layout_name or None

        raw_slides.append(
            RawSlide(
                slide_index=idx,
                title_text=title_text,
                all_text=all_text,
                text_blocks=dedup,
                image_count=img_count,
                table_count=tbl_count,
                notes=notes,
                layout_info=layout_info,
            )
        )

    return raw_slides
