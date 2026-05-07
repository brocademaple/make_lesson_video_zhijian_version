"""从 PPTX 各页导出「图片形状」内嵌位图（非整页渲染图）。"""

from __future__ import annotations

import logging
import re
import shutil
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape

logger = logging.getLogger(__name__)

_ALLOWED_EXT = frozenset(
    {"png", "jpg", "jpeg", "gif", "bmp", "tif", "tiff", "webp", "emf", "wmf"}
)


def _iter_shapes_recursive(shape: BaseShape):
    yield shape
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:  # type: ignore[attr-defined]
            yield from _iter_shapes_recursive(child)


def _normalize_ext(raw: str | None) -> str:
    if not raw:
        return "png"
    e = raw.lower().lstrip(".")
    if e not in _ALLOWED_EXT:
        return "png"
    return e


def export_picture_shapes_on_slide(
    prs: Presentation,
    slide_index: int,
    shapes_out_dir: Path,
) -> list[str]:
    """
    将指定页上所有 PICTURE 形状的内嵌图像写入 ``shapes_out_dir``。
    文件名 ``shape-0000.{ext}`` 按遍历顺序递增。
    """
    shapes_out_dir.mkdir(parents=True, exist_ok=True)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return []
    slide = prs.slides[slide_index]
    written: list[str] = []
    idx = 0
    for shape in slide.shapes:
        for s in _iter_shapes_recursive(shape):
            if s.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                img = s.image
                ext = _normalize_ext(getattr(img, "ext", None))
                blob = img.blob
                if not blob:
                    continue
                fn = f"shape-{idx:04d}.{ext}"
                (shapes_out_dir / fn).write_bytes(blob)
                written.append(fn)
                idx += 1
            except Exception as e:
                logger.warning(
                    "跳过无法导出的图片形状（页 %s）：%s",
                    slide_index,
                    e,
                )
    return written


def populate_slide_preview_folders(
    pptx_path: Path,
    previews_root: Path,
    slide_count: int,
    pngs: list[Path] | None,
) -> list[dict[str, Any]]:
    """
    在 ``previews_root`` 下为每页创建 ``slide-NNNN/``：

    - ``full.png``：与顶层 ``slide-NNNN.png`` 相同的整页渲染图（若该页有 PNG）；
    - ``shapes/shape-XXXX.ext``：该页内图片形状内嵌图。

    返回 ``shape_image_manifest`` 用条目列表。
    """
    pngs = pngs or []
    previews_root.mkdir(parents=True, exist_ok=True)
    try:
        prs = Presentation(str(pptx_path))
    except Exception:
        logger.exception("无法打开 PPTX 以导出页内图片：%s", pptx_path)
        return []

    manifest: list[dict[str, Any]] = []
    n_pres = len(prs.slides)

    for i in range(slide_count):
        slide_dir = previews_root / f"slide-{i:04d}"
        slide_dir.mkdir(parents=True, exist_ok=True)
        entry: dict[str, Any] = {
            "slide_index": i,
            "full": False,
            "shapes": [],
        }
        if i < len(pngs) and pngs[i].is_file():
            try:
                shutil.copy2(pngs[i], slide_dir / "full.png")
                entry["full"] = True
            except OSError as e:
                logger.warning("复制整页 full.png 失败 页 %s：%s", i, e)

        shapes_dir = slide_dir / "shapes"
        if i < n_pres:
            names = export_picture_shapes_on_slide(prs, i, shapes_dir)
        else:
            names = []
        entry["shapes"] = names
        manifest.append(entry)

    return manifest


def natural_shape_sort_key(p: Path) -> tuple[int, str]:
    m = re.search(r"shape-(\d+)", p.name, re.I)
    return (int(m.group(1)) if m else 0, p.name)


def list_slide_shape_files(previews_root: Path, slide_index: int) -> list[Path]:
    d = previews_root / f"slide-{slide_index:04d}" / "shapes"
    if not d.is_dir():
        return []
    files = [p for p in d.iterdir() if p.is_file()]
    return sorted(files, key=natural_shape_sort_key)
