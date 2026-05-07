from __future__ import annotations

import re
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.slide import Slide


def _iter_shapes_recursive(shape: BaseShape) -> Iterable[BaseShape]:
    yield shape
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:  # type: ignore[attr-defined]
            yield from _iter_shapes_recursive(child)


def _shape_text(shape: BaseShape) -> str:
    if not shape.has_text_frame:
        return ""
    parts: list[str] = []
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs).strip()
        if line:
            parts.append(line)
    return "\n".join(parts)


def slide_plain_text(slide: Slide) -> str:
    """提取幻灯片内所有文本形状中的可见文本（按形状遍历顺序）。"""
    chunks: list[str] = []
    for shape in slide.shapes:
        for s in _iter_shapes_recursive(shape):
            t = _shape_text(s)
            if t:
                chunks.append(t)
    # 去重连续重复（复制粘贴常见）
    deduped: list[str] = []
    for c in chunks:
        if deduped and deduped[-1] == c:
            continue
        deduped.append(c)
    return "\n".join(deduped)


def guess_slide_title(slide: Slide, body_text: str) -> str:
    """优先取标题占位符；否则取正文首行。"""
    if slide.shapes.title and slide.shapes.title.text:
        t = slide.shapes.title.text.strip()
        if t:
            return _collapse_ws(t)
    first = body_text.split("\n", 1)[0].strip() if body_text else ""
    return _collapse_ws(first)[:80] if first else "（无标题）"


def _collapse_ws(s: str) -> str:
    return re.sub(r"\s+", " ", s).strip()


def load_presentation(path: str) -> Presentation:
    return Presentation(path)


def _slide_notes(slide: Slide) -> str | None:
    try:
        if getattr(slide, "has_notes_slide", False):
            t = slide.notes_slide.notes_text_frame.text.strip()
            return t or None
    except Exception:
        pass
    return None


def _count_pictures(shape: BaseShape) -> int:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return 1
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return sum(_count_pictures(c) for c in shape.shapes)  # type: ignore[attr-defined]
    return 0


def _count_tables(shape: BaseShape) -> int:
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        return 1
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return sum(_count_tables(c) for c in shape.shapes)  # type: ignore[attr-defined]
    return 0


def slide_text_blocks(slide: Slide) -> list[str]:
    """去重后的文本块列表（用于预览与结构化展示）。"""
    chunks: list[str] = []
    for shape in slide.shapes:
        for s in _iter_shapes_recursive(shape):
            t = _shape_text(s)
            if t:
                chunks.append(t)
    deduped: list[str] = []
    for c in chunks:
        if deduped and deduped[-1] == c:
            continue
        deduped.append(c)
    return deduped


def slide_snapshot(slide: Slide, index: int) -> dict:
    """单页解析结果（供 Web JSON）。"""
    blocks = slide_text_blocks(slide)
    all_text = "\n\n".join(blocks)
    title = guess_slide_title(slide, all_text)
    img = tbl = 0
    for shape in slide.shapes:
        img += _count_pictures(shape)
        tbl += _count_tables(shape)
    layout_name = ""
    try:
        layout_name = slide.slide_layout.name or ""
    except Exception:
        pass
    return {
        "index": index,
        "title": title,
        "text": all_text,
        "text_blocks": blocks,
        "notes": _slide_notes(slide),
        "image_count": img,
        "table_count": tbl,
        "layout": layout_name or None,
    }


def parse_pptx_bytes(raw: bytes) -> list[dict]:
    """从内存中的 pptx 字节解析全部页面。"""
    import os
    import tempfile

    fd, path = tempfile.mkstemp(suffix=".pptx")
    try:
        os.write(fd, raw)
        os.close(fd)
        prs = Presentation(path)
        return [slide_snapshot(s, i) for i, s in enumerate(prs.slides)]
    finally:
        try:
            os.unlink(path)
        except OSError:
            pass
