from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ppt_course.models import CourseSlidePlan

_MARGIN = Inches(0.45)
_TITLE_TOP = Inches(0.35)
_TITLE_H = Inches(0.75)
_BODY_WIDTH = Inches(12.4)
_PALETTE = {
    "accent": RGBColor(0x1A, 0x56, 0x9E),
    "body": RGBColor(0x22, 0x22, 0x22),
    "label": RGBColor(0x0F, 0x4C, 0x81),
}


def _add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(_MARGIN, _TITLE_TOP, _BODY_WIDTH, _TITLE_H)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _PALETTE["accent"]
    p.alignment = PP_ALIGN.LEFT


def _add_section(slide, top_in: float, label: str, body: str) -> float:
    """返回下一节的 top（英寸数值）。"""
    lab_h = Inches(0.28)
    box = slide.shapes.add_textbox(_MARGIN, Inches(top_in), _BODY_WIDTH, lab_h)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = _PALETTE["label"]

    body_top = top_in + 0.32
    # 估算高度：粗略按行数
    lines = max(2, body.count("\n") + 1 + len(body) // 42)
    est_h = min(4.2, 0.28 + lines * 0.22)
    body_box = slide.shapes.add_textbox(
        _MARGIN + Inches(0.15), Inches(body_top), _BODY_WIDTH - Inches(0.15), Inches(est_h)
    )
    btf = body_box.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.text = body
    bp.font.size = Pt(15)
    bp.font.color.rgb = _PALETTE["body"]
    bp.line_spacing = 1.15

    return body_top + est_h + 0.12


def build_course_presentation(plans: list[CourseSlidePlan]) -> Presentation:
    """由课程页规划生成新版式 PPTX（空白版式 + 文本框）。"""
    prs = Presentation()
    # 使用空白布局（索引因模板而异；python-pptx 默认模板 layout 6 常为空白）
    blank_layout = prs.slide_layouts[6]

    for plan in plans:
        slide = prs.slides.add_slide(blank_layout)
        _add_title(slide, plan.slide_title)

        cursor = 1.25  # inches from top for content

        if not plan.sections:
            slide.shapes.add_textbox(
                _MARGIN, Inches(cursor), _BODY_WIDTH, Inches(1.5)
            ).text_frame.paragraphs[0].text = "（无结构化内容）"

        for label, body in plan.sections.items():
            # 单节过高时压缩字号通过 word_wrap；过高则截断提示
            if len(body) > 3500:
                body = body[:3400] + "\n…（内容过长已截断，详见备注）"
            cursor = _add_section(slide, cursor, label, body)

        if plan.notes:
            try:
                tf = slide.notes_slide.notes_text_frame
                tf.text = "【原始页摘录 / 口播参考】\n" + plan.notes[:12000]
            except (AttributeError, ValueError):
                pass

    return prs


def save_presentation(prs: Presentation, path: str) -> None:
    prs.save(path)
